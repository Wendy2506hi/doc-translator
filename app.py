# 📄 文档翻译小工具（支持 Word、PPT 和 PDF，PDF保留原排版）

# 功能：
# 1. 网页上传 .docx / .pptx / .pdf 文件
# 2. 自动提取文本并调用 ChatGPT 翻译为目标语言
# 3. 保留 Word/PPT 原排版；PDF 原位替换文字后重新生成 PDF
# 4. 下载翻译后的文件（同格式输出）

# 环境要求：
# pip install openai python-docx python-pptx streamlit PyMuPDF

import os
import openai
import streamlit as st
from docx import Document
from pptx import Presentation
from tempfile import NamedTemporaryFile
import fitz  # PyMuPDF，用于PDF处理

# 替换为你的 OpenAI API Key
openai.api_key = os.getenv("OPENAI_API_KEY")

def translate_text(text, target_language="中文"):
    if not text.strip():
        return text
    response = openai.ChatCompletion.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": f"You are a professional translator. Translate all content into {target_language}, preserving formatting and meaning."},
            {"role": "user", "content": text},
        ],
        temperature=0.3
    )
    return response["choices"][0]["message"]["content"]

def translate_word(file, target_language):
    doc = Document(file)
    for para in doc.paragraphs:
        para.text = translate_text(para.text, target_language)
    temp_file = NamedTemporaryFile(delete=False, suffix=".docx")
    doc.save(temp_file.name)
    return temp_file.name

def translate_ppt(file, target_language):
    ppt = Presentation(file)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape.text = translate_text(shape.text, target_language)
    temp_file = NamedTemporaryFile(delete=False, suffix=".pptx")
    ppt.save(temp_file.name)
    return temp_file.name

def translate_pdf_preserve_layout(file, target_language):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    for page in doc:
        blocks = page.get_text("dict")["blocks"]
        for b in blocks:
            if b["type"] == 0:  # text block
                for line in b["lines"]:
                    for span in line["spans"]:
                        orig_text = span["text"]
                        if orig_text.strip():
                            translated = translate_text(orig_text, target_language)
                            # 覆盖原文本
                            page.insert_text(
                                (span["bbox"][0], span["bbox"][1]),
                                translated,
                                fontname=span["font"] if span["font"] else "helv",
                                fontsize=span["size"],
                                color=(0, 0, 0),
                                overlay=True
                            )
    temp_file = NamedTemporaryFile(delete=False, suffix=".pdf")
    doc.save(temp_file.name)
    return temp_file.name

# Streamlit UI
st.set_page_config(page_title="📄 文档翻译助手", layout="centered")
st.title("📄 文档翻译助手（ChatGPT驱动）")
st.markdown("上传 Word (.docx)、PowerPoint (.pptx) 或 PDF (.pdf) 文件，选择目标语言，即可自动翻译并下载结果。\n\n✅ PDF 保留版式、直接在原位替换翻译文字。")

uploaded_file = st.file_uploader("上传文件", type=["docx", "pptx", "pdf"])
target_language = st.selectbox("选择目标语言", ["中文", "英文", "德文", "日文"])

if uploaded_file and st.button("开始翻译"):
    with st.spinner("翻译中，请稍候..."):
        suffix = uploaded_file.name.lower().split(".")[-1]
        if suffix == "docx":
            translated_path = translate_word(uploaded_file, target_language)
            mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        elif suffix == "pptx":
            translated_path = translate_ppt(uploaded_file, target_language)
            mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif suffix == "pdf":
            translated_path = translate_pdf_preserve_layout(uploaded_file, target_language)
            mime_type = "application/pdf"
        else:
            st.error("仅支持 .docx / .pptx / .pdf 文件")
            st.stop()

    with open(translated_path, "rb") as f:
        st.success("🎉 翻译完成！")
        st.download_button(
            label="📥 下载翻译文件",
            data=f.read(),
            file_name=f"translated_{uploaded_file.name}",
            mime=mime_type
        )
